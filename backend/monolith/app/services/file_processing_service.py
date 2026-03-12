import csv
import io
import os
import re

import fitz  # PyMuPDF

from app.config import PROCESSED_DIR

os.makedirs(PROCESSED_DIR, exist_ok=True)


def _clean_and_normalize_text(raw_text: str) -> str:
    text = raw_text.lower()
    text = text.replace("\n", " ")
    text = re.sub(r"\s+", " ", text).strip()
    return text


# ---- Extractors por formato ----

def _extract_pdf(path: str) -> str:
    try:
        doc = fitz.open(path)
        text = "".join(page.get_text() for page in doc)
        doc.close()
    except Exception as e:
        raise ValueError(f"PyMuPDF no pudo procesar el archivo: {e}")

    if len(text.strip()) < 50:
        try:
            from app.services import s3_service
            text = s3_service.extract_text_from_pdf_images_via_textract(path)
        except Exception:
            pass

    return text


def _extract_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_csv(path: str) -> str:
    rows = []
    with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(rows)


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Hoja: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_xls(path: str) -> str:
    import xlrd
    wb = xlrd.open_workbook(path)
    parts = []
    for sheet in wb.sheets():
        parts.append(f"[Hoja: {sheet.name}]")
        for rx in range(sheet.nrows):
            row = sheet.row(rx)
            cells = [str(c.value).strip() for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Diapositiva {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


_EXTRACTORS = {
    "pdf": _extract_pdf,
    "txt": _extract_txt,
    "csv": _extract_csv,
    "docx": _extract_docx,
    "xlsx": _extract_xlsx,
    "xls": _extract_xls,
    "pptx": _extract_pptx,
}


def process_file_to_text(local_path: str) -> str:
    """Extrae texto del archivo según su extensión y guarda un .txt normalizado."""
    _, ext = os.path.splitext(local_path)
    ext = ext.lower().lstrip(".")

    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(f"Formato no soportado: .{ext}")

    raw_text = extractor(local_path)
    cleaned = _clean_and_normalize_text(raw_text)

    base = os.path.basename(local_path)
    out_path = os.path.join(PROCESSED_DIR, os.path.splitext(base)[0] + ".txt")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(cleaned)
    return out_path
